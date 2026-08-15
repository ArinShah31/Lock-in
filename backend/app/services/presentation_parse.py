"""Extract per-slide text, pictures, and diagram regions from a PPTX."""

from __future__ import annotations

import re
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

_DIAGRAM_WORDS = (
    "flowchart",
    "flow chart",
    "diagram",
    "architecture",
    "pipeline",
    "workflow",
    "state machine",
    "er diagram",
    "uml",
)

_BOILERPLATE_RE = re.compile(
    r"(university|institute|school of computing|ay\s*20\d{2}|"
    r"sem-?\s*ii|class\s*-\s*s\.y|<division>|mit art|technology university|"
    r"\bprof\.|\bprofessor\b|academic year)",
    re.I,
)


def _emu_pct(value, total) -> float:
    if not total:
        return 0.0
    try:
        return max(0.0, min(1.0, float(value) / float(total)))
    except Exception:
        return 0.0


def _shape_box(shape, slide_w, slide_h) -> dict[str, float]:
    left = getattr(shape, "left", 0) or 0
    top = getattr(shape, "top", 0) or 0
    width = getattr(shape, "width", 0) or 0
    height = getattr(shape, "height", 0) or 0
    return {
        "x": round(_emu_pct(left, slide_w), 4),
        "y": round(_emu_pct(top, slide_h), 4),
        "w": round(max(_emu_pct(width, slide_w), 0.04), 4),
        "h": round(max(_emu_pct(height, slide_h), 0.04), 4),
    }


def _shape_kind(shape) -> str:
    try:
        shape_type = shape.shape_type
    except Exception:
        return "shape"
    picture_types = {MSO_SHAPE_TYPE.PICTURE}
    if hasattr(MSO_SHAPE_TYPE, "LINKED_PICTURE"):
        picture_types.add(MSO_SHAPE_TYPE.LINKED_PICTURE)
    if shape_type in picture_types:
        return "picture"
    if shape_type == MSO_SHAPE_TYPE.CHART:
        return "chart"
    if shape_type == MSO_SHAPE_TYPE.TABLE:
        return "table"
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if shape_type in {MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM}:
        return "connector"
    if getattr(shape, "has_text_frame", False) and _shape_text(shape):
        return "text"
    return "shape"


def _shape_text(shape) -> str:
    if getattr(shape, "has_text_frame", False):
        return (shape.text or "").strip()
    if hasattr(shape, "text"):
        try:
            return (shape.text or "").strip()
        except Exception:
            return ""
    return ""


def _element_role(kind: str, box: dict) -> str:
    y = float(box.get("y") or 0)
    w = float(box.get("w") or 0)
    h = float(box.get("h") or 0)
    area = w * h
    if kind in {"picture", "shape", "group"} and (area >= 0.7 or (w >= 0.92 and h >= 0.88)):
        return "background"
    if y <= 0.1 and h <= 0.18:
        return "chrome"
    if y + h >= 0.9 and h <= 0.16:
        return "chrome"
    if kind == "picture" and area <= 0.08 and y <= 0.2:
        return "chrome"
    return "content"


def _element(kind: str, collected: list[dict], box: dict, text: str = "", name: str = "") -> dict:
    index = len(collected)
    return {
        "index": index,
        "element_id": f"{kind}_{index}",
        "kind": kind,
        "role": _element_role(kind, box),
        "name": name,
        "text": text,
        **box,
    }


def _shape_name(shape) -> str:
    try:
        return str(getattr(shape, "name", "") or "")
    except Exception:
        return ""


def _collect_shapes(shapes, slide_w, slide_h, collected: list[dict], texts: list[str], stats: dict) -> None:
    for shape in shapes:
        kind = _shape_kind(shape)
        box = _shape_box(shape, slide_w, slide_h)
        name = _shape_name(shape)
        if kind == "group":
            role = _element_role("group", box)
            if role != "content":
                stats["chrome"] += 1
            else:
                stats["visuals"] += 1
            collected.append(_element("group", collected, box, name=name))
            try:
                _collect_shapes(shape.shapes, slide_w, slide_h, collected, texts, stats)
            except Exception:
                pass
            continue
        if kind in {"picture", "chart", "table"}:
            role = _element_role(kind, box)
            if role != "content":
                stats["chrome"] += 1
            else:
                stats["visuals"] += 1
            collected.append(_element(kind, collected, box, text=_shape_text(shape), name=name))
            continue
        if kind == "connector":
            stats["connectors"] += 1
            continue
        text = _shape_text(shape)
        if text:
            texts.append(text)
            collected.append(_element("text", collected, box, text=text, name=name))
            continue
        collected.append(_element("shape", collected, box, name=name))


def shape_role(shape: dict | None) -> str:
    if not shape:
        return "content"
    role = str(shape.get("role") or "")
    if role in {"background", "chrome", "content"}:
        return role
    return _element_role(str(shape.get("kind") or "shape"), shape)


def line_is_boilerplate(line: str) -> bool:
    raw = re.sub(r"[\s\u000b]+", " ", line or "").strip()
    if not raw:
        return True
    return bool(len(raw) < 120 and _BOILERPLATE_RE.search(raw))


def content_text(extracted: str, repeated: set[str] | None = None) -> str:
    lines: list[str] = []
    for ln in (extracted or "").splitlines():
        norm = re.sub(r"[\s\u000b]+", " ", ln).strip()
        if not norm or line_is_boilerplate(norm):
            continue
        if repeated and norm.lower() in repeated:
            continue
        lines.append(norm)
    return "\n".join(lines).strip()


def shape_is_mostly_chrome_text(shape: dict | None) -> bool:
    if not shape:
        return False
    raw = str(shape.get("text") or "")
    if not raw.strip():
        return False
    cleaned = content_text(raw)
    if not cleaned:
        return True
    area = float(shape.get("w") or 0) * float(shape.get("h") or 0)
    return len(cleaned) / max(len(raw), 1) < 0.45 and area > 0.2


def is_content_visual(shape: dict | None) -> bool:
    if not shape or shape_role(shape) != "content":
        return False
    return str(shape.get("kind") or "") in {"picture", "chart", "group", "table"}


def is_highlightable(shape: dict | None) -> bool:
    if not shape or shape_role(shape) != "content":
        return False
    area = float(shape.get("w") or 0) * float(shape.get("h") or 0)
    if area >= 0.78:
        return False
    if shape_is_mostly_chrome_text(shape):
        return False
    return True


def slide_looks_like_diagram(shapes: list[dict] | None, extracted_text: str = "") -> bool:
    items = shapes or []
    visuals = [item for item in items if is_content_visual(item)]
    if any(str(item.get("kind") or "") in {"chart", "table"} for item in visuals):
        return True
    sized = []
    for item in visuals:
        area = float(item.get("w") or 0) * float(item.get("h") or 0)
        if 0.1 <= area <= 0.75:
            sized.append(item)
    if len(sized) >= 2:
        return True
    if len(sized) == 1 and str(sized[0].get("kind") or "") in {"picture", "group"}:
        return True
    body = content_text(extracted_text).lower()
    return any(word in body for word in _DIAGRAM_WORDS)


def _repeated_header_lines(slides: list[dict]) -> set[str]:
    counter: Counter[str] = Counter()
    for slide in slides:
        seen: set[str] = set()
        for ln in (slide.get("extracted_text") or "").splitlines():
            norm = re.sub(r"[\s\u000b]+", " ", ln).strip().lower()
            if len(norm) < 10 or norm in seen:
                continue
            seen.add(norm)
            counter[norm] += 1
    n = max(len(slides), 1)
    threshold = max(3, int(n * 0.35))
    return {key for key, count in counter.items() if count >= threshold}


def extract_slides(file_path: str) -> list[dict]:
    presentation = Presentation(file_path)
    slide_w = presentation.slide_width or 1
    slide_h = presentation.slide_height or 1
    slides: list[dict] = []

    for index, slide in enumerate(presentation.slides):
        collected: list[dict] = []
        texts: list[str] = []
        stats = {"visuals": 0, "connectors": 0, "chrome": 0}
        try:
            _collect_shapes(slide.shapes, slide_w, slide_h, collected, texts, stats)
        except Exception:
            pass
        extracted = "\n".join(texts).strip()
        slides.append(
            {
                "index": index,
                "extracted_text": extracted,
                "script": extracted or f"Slide {index + 1}.",
                "shapes": collected,
                "connectors": stats["connectors"],
            }
        )

    repeated = _repeated_header_lines(slides)
    for slide in slides:
        original = slide["extracted_text"]
        cleaned = content_text(original, repeated)
        if cleaned:
            slide["extracted_text"] = cleaned
        if slide["script"] == original:
            slide["script"] = slide["extracted_text"] or f"Slide {slide['index'] + 1}."
        slide["has_diagram"] = int(slide.pop("connectors", 0) or 0) >= 3 or slide_looks_like_diagram(
            slide["shapes"], slide["extracted_text"]
        )

    if not slides:
        slides.append(
            {
                "index": 0,
                "extracted_text": Path(file_path).stem,
                "script": f"This presentation is titled {Path(file_path).stem}.",
                "shapes": [],
                "has_diagram": False,
            }
        )
    return slides
