from pptx import Presentation


def extract_text(file_path: str) -> tuple[str, int]:
    presentation = Presentation(file_path)

    text = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                value = shape.text.strip()

                if value:
                    text.append(value)

    return "\n".join(text), len(presentation.slides)